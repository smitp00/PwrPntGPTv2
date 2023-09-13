from django.http import HttpResponse
from django.shortcuts import render, redirect
from pptx import Presentation
from docx import Document
import boto3
import tempfile  # Added for temporary file handling
from .gpt_service import get_summary_from_gpt  # Place gpt_service.py in the my_app directory
from decouple import config


# Initialize S3 client
s3_client = boto3.client(
    's3',
    aws_access_key_id= config('AWS_ACCESS_KEY_ID'), #aws_access_key_id=os.environ.get('AWS_ACCESS_KEY_ID',''),  # Use environment variables
    aws_secret_access_key=config('AWS_SECRET_ACCESS_KEY'),#aws_secret_access_key=os.environ.get('AWS_SECRET_ACCESS_KEY',''),
    region_name='ca-central-1'
)


def upload_file_to_s3(file, filename):
    s3_client.upload_fileobj(file, 'pwrpntgptbucket', filename)

def download_file_from_s3(filename):
    with tempfile.NamedTemporaryFile(delete=False) as f:
        s3_client.download_file('pwrpntgptbucket', filename, f.name)
        return f.name

#Activates from post request and user-inputted ppt or word files
def upload(request):
    if request.method == 'POST':
        ppt_file = request.FILES.get('ppt_file')
        word_file = request.FILES.get('word_file')

        if ppt_file and ppt_file.name.endswith('.pptx'):
            try:
                upload_file_to_s3(ppt_file, ppt_file.name)

                temp_file_path_ppt = download_file_from_s3(ppt_file.name)
                prs = Presentation(temp_file_path_ppt)
                slides_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slides_text.append(run.text)

                        # Handle table shape
                        if shape.shape_type == 19:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    slides_text.append(cell.text)

                if not slides_text:
                    return render(request, 'error.html', {'error_message': 'The presentation appears to be empty.'})
                
                summarized_content = get_summary_from_gpt(' '.join(slides_text))
                return render(request, 'summaryppt.html', {'original_content': slides_text, 'summarized_content': summarized_content})
                
            except Exception as e:
                return render(request, 'error.html', {'error_message': str(e)})

        elif word_file and word_file.name.endswith('.docx'):
            upload_file_to_s3(word_file, word_file.name)
            temp_file_path_word = download_file_from_s3(word_file.name)

            document = Document(temp_file_path_word)
            doc_text = []
            for para in document.paragraphs:
                doc_text.append(para.text)
            
            # Combine all the text into a single string
            full_text = '\n'.join(doc_text)
            # Get the summarized content
            summary = get_summary_from_gpt(full_text)
            return render(request, 'summaryword.html', {'summary': summary})
            
    return render(request, 'upload.html')